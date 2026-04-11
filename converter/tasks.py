import os
import tempfile
from pathlib import Path
from celery import shared_task
from django.utils import timezone


@shared_task(bind=True)
def compress_pdf_task(self, job_id, level='recommended'):
    from converter.models import ConversionJob
    from converter.utils import get_output_path
    job = ConversionJob.objects.get(id=job_id)
    try:
        job.status = 'processing'; job.save()
        import fitz
        abs_path, rel_path = get_output_path(job.input_file.name, '.pdf')
        doc = fitz.open(job.input_file.path)
        if level == 'extreme':
            save_kwargs = {'garbage': 4, 'deflate': True, 'deflate_images': True,
                           'deflate_fonts': True, 'clean': True}
        elif level == 'less':
            save_kwargs = {'garbage': 1, 'deflate': True, 'clean': False}
        else:  # recommended
            save_kwargs = {'garbage': 4, 'deflate': True, 'clean': True}
        doc.save(abs_path, **save_kwargs)
        doc.close()
        job.output_file = rel_path
        job.output_size = os.path.getsize(abs_path)
        job.status = 'done'; job.save()
    except Exception as e:
        job.status = 'failed'; job.error_message = str(e); job.save()
        raise


# ════════════════════════════════════════════════════════════════
#  REPLACEMENT 2 — converter/tasks.py
#  Replace the entire merge_pdfs_task function
# ════════════════════════════════════════════════════════════════
 
@shared_task(bind=True)
def merge_pdfs_task(self, job_id, extra_paths=None):
    from converter.models import ConversionJob
    from converter.utils import get_output_path
    job = ConversionJob.objects.get(id=job_id)
    try:
        job.status = 'processing'; job.save()
        import fitz
 
        # First file is the main uploaded file
        input_files = [job.input_file.path]
 
        # Additional files passed as task kwargs (no longer via error_message)
        if extra_paths:
            input_files.extend([p for p in extra_paths if os.path.exists(p)])
 
        abs_path, rel_path = get_output_path(job.input_file.name, '.pdf')
        merged = fitz.open()
        for fp in input_files:
            doc = fitz.open(fp)
            merged.insert_pdf(doc)
            doc.close()
        merged.save(abs_path, garbage=4, deflate=True)
        merged.close()
 
        # Clean up the extra files from disk
        for fp in (extra_paths or []):
            try:
                os.remove(fp)
            except Exception:
                pass
 
        job.output_file = rel_path
        job.output_size = os.path.getsize(abs_path)
        job.status = 'done'; job.save()
    except Exception as e:
        job.status = 'failed'; job.error_message = str(e); job.save()
        raise


@shared_task(bind=True)
def split_pdf_task(self, job_id, start_page=None, end_page=None):
    from converter.models import ConversionJob
    from converter.utils import create_zip
    job = ConversionJob.objects.get(id=job_id)
    try:
        job.status = 'processing'; job.save()
        import fitz
        doc = fitz.open(job.input_file.path)
        total = len(doc)
        s = max(0, (int(start_page) - 1) if start_page else 0)
        e = min(total, int(end_page) if end_page else total)
        e = max(s + 1, e)
        tmp_dir = tempfile.mkdtemp()
        files_dict = {}
        for i in range(s, e):
            out = fitz.open(); out.insert_pdf(doc, from_page=i, to_page=i)
            fp = os.path.join(tmp_dir, f"page_{i+1}.pdf")
            out.save(fp); out.close()
            files_dict[f"page_{i+1}.pdf"] = fp
        doc.close()
        zip_abs, zip_rel = create_zip(files_dict, 'split_pages.zip')
        for fp in files_dict.values():
            try: os.remove(fp)
            except: pass
        job.output_file = zip_rel
        job.output_size = os.path.getsize(zip_abs)
        job.status = 'done'; job.save()
    except Exception as e:
        job.status = 'failed'; job.error_message = str(e); job.save()
        raise


@shared_task(bind=True)
def pdf_to_images_task(self, job_id, dpi=150, img_format='png'):
    from converter.models import ConversionJob
    from converter.utils import create_zip
    job = ConversionJob.objects.get(id=job_id)
    try:
        job.status = 'processing'; job.save()
        import fitz
        from PIL import Image
        import io

        dpi = int(dpi) if dpi else 150
        fmt = (img_format or 'png').lower()
        if fmt not in ('png', 'jpeg', 'webp'):
            fmt = 'png'

        ext_map  = {'png': '.png', 'jpeg': '.jpg', 'webp': '.webp'}
        pil_fmt  = {'png': 'PNG',  'jpeg': 'JPEG', 'webp': 'WEBP'}
        file_ext = ext_map[fmt]

        doc = fitz.open(job.input_file.path)
        tmp_dir = tempfile.mkdtemp()
        files_dict = {}
        mat = fitz.Matrix(dpi / 72, dpi / 72)

        for i, page in enumerate(doc):
            pix = page.get_pixmap(matrix=mat, alpha=False)
            img_filename = f"page_{i + 1}{file_ext}"
            img_abs = os.path.join(tmp_dir, img_filename)

            if fmt == 'png':
                pix.save(img_abs)
            else:
                img_bytes = pix.tobytes('png')
                img = Image.open(io.BytesIO(img_bytes))
                if img.mode in ('RGBA', 'P', 'LA'):
                    bg = Image.new('RGB', img.size, (255, 255, 255))
                    if img.mode == 'P':
                        img = img.convert('RGBA')
                    bg.paste(img, mask=img.split()[3] if img.mode == 'RGBA' else None)
                    img = bg
                elif img.mode != 'RGB':
                    img = img.convert('RGB')
                img.save(img_abs, pil_fmt[fmt], quality=85, optimize=True)

            files_dict[img_filename] = img_abs

        doc.close()
        zip_abs, zip_rel = create_zip(files_dict, f'pdf_images_{fmt}.zip')
        for fp in files_dict.values():
            try: os.remove(fp)
            except: pass
        job.output_file = zip_rel
        job.output_size = os.path.getsize(zip_abs)
        job.status = 'done'; job.save()
    except Exception as e:
        job.status = 'failed'; job.error_message = str(e); job.save()
        raise


@shared_task(bind=True)
def docx_to_pdf_task(self, job_id):
    from converter.models import ConversionJob
    from converter.utils import get_output_path
    job = ConversionJob.objects.get(id=job_id)
    try:
        job.status = 'processing'; job.save()
        from docx import Document
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer

        abs_path, rel_path = get_output_path(job.input_file.name, '.pdf')
        doc_word = Document(job.input_file.path)
        pdf_doc = SimpleDocTemplate(
            abs_path, pagesize=A4,
            rightMargin=inch, leftMargin=inch,
            topMargin=inch, bottomMargin=inch
        )
        styles = getSampleStyleSheet()
        normal = ParagraphStyle('N', parent=styles['Normal'], fontSize=11, leading=16, spaceAfter=8)
        h1 = ParagraphStyle('H1', parent=styles['Heading1'], fontSize=15, leading=22, spaceAfter=12, spaceBefore=16)
        h2 = ParagraphStyle('H2', parent=styles['Heading2'], fontSize=13, leading=18, spaceAfter=10, spaceBefore=12)

        story = []
        for para in doc_word.paragraphs:
            raw = para.text
            if not raw or not raw.strip():
                story.append(Spacer(1, 0.12 * inch)); continue
            text = raw.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            sname = (para.style.name if para.style and para.style.name else '')
            if 'Heading 1' in sname or 'Title' in sname:
                story.append(Paragraph(text, h1))
            elif 'Heading 2' in sname or 'Subtitle' in sname:
                story.append(Paragraph(text, h2))
            else:
                story.append(Paragraph(text, normal))
        if not story:
            story.append(Paragraph("(Empty document)", normal))
        pdf_doc.build(story)
        job.output_file = rel_path
        job.output_size = os.path.getsize(abs_path)
        job.status = 'done'; job.save()
    except Exception as e:
        job.status = 'failed'; job.error_message = str(e); job.save()
        raise


@shared_task(bind=True)
def txt_to_pdf_task(self, job_id):
    from converter.models import ConversionJob
    from converter.utils import get_output_path
    job = ConversionJob.objects.get(id=job_id)
    try:
        job.status = 'processing'; job.save()
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import ParagraphStyle
        from reportlab.lib.units import inch
        from reportlab.platypus import SimpleDocTemplate, Paragraph

        abs_path, rel_path = get_output_path(job.input_file.name, '.pdf')
        with open(job.input_file.path, 'r', encoding='utf-8', errors='replace') as f:
            content = f.read()
        pdf_doc = SimpleDocTemplate(abs_path, pagesize=A4,
                                    rightMargin=inch, leftMargin=inch,
                                    topMargin=inch, bottomMargin=inch)
        code_style = ParagraphStyle('Code', fontName='Courier', fontSize=10, leading=14)
        story = []
        for line in content.split('\n'):
            safe = line.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;') or '&nbsp;'
            story.append(Paragraph(safe, code_style))
        pdf_doc.build(story)
        job.output_file = rel_path
        job.output_size = os.path.getsize(abs_path)
        job.status = 'done'; job.save()
    except Exception as e:
        job.status = 'failed'; job.error_message = str(e); job.save()
        raise


@shared_task(bind=True)
def img_to_pdf_task(self, job_id):
    from converter.models import ConversionJob
    from converter.utils import get_output_path
    job = ConversionJob.objects.get(id=job_id)
    try:
        job.status = 'processing'; job.save()
        from PIL import Image
        from reportlab.lib.pagesizes import A4
        from reportlab.platypus import SimpleDocTemplate, Image as RLImage
        from reportlab.lib.units import inch

        abs_path, rel_path = get_output_path(job.input_file.name, '.pdf')
        img = Image.open(job.input_file.path)
        if img.mode in ('RGBA', 'P', 'LA'):
            img = img.convert('RGB')
        elif img.mode != 'RGB':
            img = img.convert('RGB')
        tmp_jpg = abs_path + '_tmp.jpg'
        img.save(tmp_jpg, 'JPEG', quality=95)
        w, h = img.size
        pw, ph = A4
        ratio = min((pw - 2 * inch) / w, (ph - 2 * inch) / h)
        pdf_doc = SimpleDocTemplate(abs_path, pagesize=A4,
                                    rightMargin=inch, leftMargin=inch,
                                    topMargin=inch, bottomMargin=inch)
        pdf_doc.build([RLImage(tmp_jpg, width=w * ratio, height=h * ratio)])
        try: os.remove(tmp_jpg)
        except: pass
        job.output_file = rel_path
        job.output_size = os.path.getsize(abs_path)
        job.status = 'done'; job.save()
    except Exception as e:
        job.status = 'failed'; job.error_message = str(e); job.save()
        raise


@shared_task(bind=True)
def jpg_to_png_task(self, job_id):
    from converter.models import ConversionJob
    from converter.utils import get_output_path
    job = ConversionJob.objects.get(id=job_id)
    try:
        job.status = 'processing'; job.save()
        from PIL import Image
        abs_path, rel_path = get_output_path(job.input_file.name, '.png')
        img = Image.open(job.input_file.path)
        img.load()
        img.save(abs_path, 'PNG')
        job.output_file = rel_path
        job.output_size = os.path.getsize(abs_path)
        job.status = 'done'; job.save()
    except Exception as e:
        job.status = 'failed'; job.error_message = str(e); job.save()
        raise


@shared_task(bind=True)
def png_to_jpg_task(self, job_id, quality=85):
    from converter.models import ConversionJob
    from converter.utils import get_output_path
    job = ConversionJob.objects.get(id=job_id)
    try:
        job.status = 'processing'; job.save()
        from PIL import Image
        quality = int(quality) if quality else 85
        quality = max(10, min(95, quality))
        abs_path, rel_path = get_output_path(job.input_file.name, '.jpg')
        img = Image.open(job.input_file.path)
        img.load()
        if img.mode in ('RGBA', 'P', 'LA'):
            bg = Image.new('RGB', img.size, (255, 255, 255))
            if img.mode == 'P': img = img.convert('RGBA')
            bg.paste(img, mask=img.split()[3] if img.mode == 'RGBA' else None)
            img = bg
        elif img.mode != 'RGB':
            img = img.convert('RGB')
        img.save(abs_path, 'JPEG', quality=quality, optimize=True)
        job.output_file = rel_path
        job.output_size = os.path.getsize(abs_path)
        job.status = 'done'; job.save()
    except Exception as e:
        job.status = 'failed'; job.error_message = str(e); job.save()
        raise


@shared_task(bind=True)
def resize_image_task(self, job_id, width=None, height=None):
    from converter.models import ConversionJob
    from converter.utils import get_output_path
    job = ConversionJob.objects.get(id=job_id)
    try:
        job.status = 'processing'; job.save()
        from PIL import Image
        input_path = job.input_file.path
        ext = Path(input_path).suffix.lower()
        fmt_map = {'.jpg': '.jpg', '.jpeg': '.jpg', '.png': '.png',
                   '.bmp': '.bmp', '.tiff': '.tiff', '.webp': '.webp'}
        out_ext = fmt_map.get(ext, '.png')
        abs_path, rel_path = get_output_path(job.input_file.name, out_ext)
        img = Image.open(input_path); img.load()
        ow, oh = img.size
        if width and height:   new_size = (int(width), int(height))
        elif width:            new_size = (int(width), int(oh * int(width) / ow))
        elif height:           new_size = (int(ow * int(height) / oh), int(height))
        else:                  new_size = (ow // 2, oh // 2)
        img = img.resize(new_size, Image.LANCZOS)
        fmt = {'jpg': 'JPEG', 'jpeg': 'JPEG', 'png': 'PNG',
               'bmp': 'BMP', 'tiff': 'TIFF', 'webp': 'WEBP'}.get(out_ext.lstrip('.'), 'PNG')
        if fmt == 'JPEG' and img.mode != 'RGB': img = img.convert('RGB')
        img.save(abs_path, fmt)
        job.output_file = rel_path
        job.output_size = os.path.getsize(abs_path)
        job.status = 'done'; job.save()
    except Exception as e:
        job.status = 'failed'; job.error_message = str(e); job.save()
        raise


@shared_task
def cleanup_expired_jobs():
    from converter.models import ConversionJob
    expired = ConversionJob.objects.filter(expires_at__lt=timezone.now())
    count = 0
    for job in expired:
        for f in [job.input_file, job.output_file]:
            if f:
                try:
                    if os.path.exists(f.path): os.remove(f.path)
                except: pass
        job.delete(); count += 1
    return f"Cleaned {count} expired jobs"


@shared_task(bind=True)
def pdf_to_word_task(self, job_id):
    from converter.models import ConversionJob
    from converter.utils import get_output_path
    job = ConversionJob.objects.get(id=job_id)
    try:
        job.status = 'processing'; job.save()
        import fitz
        from docx import Document
        from docx.shared import Pt

        abs_path, rel_path = get_output_path(job.input_file.name, '.docx')
        doc_pdf = fitz.open(job.input_file.path)
        doc_word = Document()
        style = doc_word.styles['Normal']
        style.font.name = 'Calibri'
        style.font.size = Pt(11)

        for page_num, page in enumerate(doc_pdf):
            if page_num > 0:
                doc_word.add_page_break()
            blocks = page.get_text("dict")["blocks"]
            for block in blocks:
                if block.get("type") != 0:
                    continue
                for line in block.get("lines", []):
                    line_text = ""
                    max_size = 0
                    is_bold = False
                    for span in line.get("spans", []):
                        line_text += span.get("text", "")
                        size = span.get("size", 11)
                        if size > max_size:
                            max_size = size
                        if span.get("flags", 0) & 16:
                            is_bold = True
                    line_text = line_text.strip()
                    if not line_text:
                        continue
                    if max_size >= 18:
                        doc_word.add_heading(line_text, level=1)
                    elif max_size >= 14:
                        doc_word.add_heading(line_text, level=2)
                    elif is_bold and max_size >= 12:
                        doc_word.add_heading(line_text, level=3)
                    else:
                        doc_word.add_paragraph(line_text)

        doc_pdf.close()
        doc_word.save(abs_path)
        job.output_file = rel_path
        job.output_size = os.path.getsize(abs_path)
        job.status = 'done'; job.save()
    except Exception as e:
        job.status = 'failed'; job.error_message = str(e); job.save()
        raise


@shared_task(bind=True)
def any_to_pdf_task(self, job_id):
    from converter.models import ConversionJob
    from converter.utils import get_output_path
    job = ConversionJob.objects.get(id=job_id)
    try:
        job.status = 'processing'; job.save()
        input_path = job.input_file.path
        ext = Path(input_path).suffix.lower()
        abs_path, rel_path = get_output_path(job.input_file.name, '.pdf')

        if ext == '.pdf':
            import fitz
            doc = fitz.open(input_path)
            doc.save(abs_path, garbage=4, deflate=True)
            doc.close()

        elif ext in ('.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.webp', '.gif'):
            from PIL import Image
            from reportlab.lib.pagesizes import A4
            from reportlab.platypus import SimpleDocTemplate, Image as RLImage
            from reportlab.lib.units import inch
            img = Image.open(input_path); img.load()
            if img.mode == 'RGBA':
                bg = Image.new('RGB', img.size, (255, 255, 255))
                bg.paste(img, mask=img.split()[3])
                img = bg
            elif img.mode != 'RGB':
                img = img.convert('RGB')
            tmp = abs_path + '_tmp.jpg'
            img.save(tmp, 'JPEG', quality=95)
            w, h = img.size
            pw, ph = A4
            ratio = min((pw - 2 * inch) / w, (ph - 2 * inch) / h)
            pdf_doc = SimpleDocTemplate(abs_path, pagesize=A4,
                                        rightMargin=inch, leftMargin=inch,
                                        topMargin=inch, bottomMargin=inch)
            pdf_doc.build([RLImage(tmp, width=w * ratio, height=h * ratio)])
            try: os.remove(tmp)
            except: pass

        elif ext == '.txt':
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import ParagraphStyle
            from reportlab.lib.units import inch
            from reportlab.platypus import SimpleDocTemplate, Paragraph
            with open(input_path, 'r', encoding='utf-8', errors='replace') as f:
                content = f.read()
            pdf_doc = SimpleDocTemplate(abs_path, pagesize=A4,
                                        rightMargin=inch, leftMargin=inch,
                                        topMargin=inch, bottomMargin=inch)
            style = ParagraphStyle('Code', fontName='Courier', fontSize=10, leading=14)
            story = [Paragraph(l.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;') or '&nbsp;', style)
                     for l in content.split('\n')]
            pdf_doc.build(story)

        elif ext in ('.docx', '.doc'):
            from docx import Document
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import inch
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            doc_word = Document(input_path)
            pdf_doc = SimpleDocTemplate(abs_path, pagesize=A4,
                                        rightMargin=inch, leftMargin=inch,
                                        topMargin=inch, bottomMargin=inch)
            styles = getSampleStyleSheet()
            normal = ParagraphStyle('N', parent=styles['Normal'], fontSize=11, leading=16, spaceAfter=8)
            h1 = ParagraphStyle('H1', parent=styles['Heading1'], fontSize=15, leading=22, spaceAfter=12)
            story = []
            for para in doc_word.paragraphs:
                raw = para.text
                if not raw.strip():
                    story.append(Spacer(1, 0.1 * inch)); continue
                text = raw.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                sname = para.style.name if para.style else ''
                story.append(Paragraph(text, h1 if 'Heading' in sname or 'Title' in sname else normal))
            if not story: story.append(Paragraph('(Empty document)', normal))
            pdf_doc.build(story)

        else:
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import ParagraphStyle
            from reportlab.lib.units import inch
            from reportlab.platypus import SimpleDocTemplate, Paragraph
            try:
                with open(input_path, 'r', encoding='utf-8', errors='replace') as f:
                    content = f.read()
            except Exception:
                content = f'[Binary file: {Path(input_path).name}]'
            pdf_doc = SimpleDocTemplate(abs_path, pagesize=A4,
                                        rightMargin=inch, leftMargin=inch,
                                        topMargin=inch, bottomMargin=inch)
            style = ParagraphStyle('Code', fontName='Courier', fontSize=10, leading=14)
            story = [Paragraph(l.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;') or '&nbsp;', style)
                     for l in content.split('\n')]
            pdf_doc.build(story)

        job.output_file = rel_path
        job.output_size = os.path.getsize(abs_path)
        job.status = 'done'; job.save()
    except Exception as e:
        job.status = 'failed'; job.error_message = str(e); job.save()
        raise


# ─────────────────────────────────────────────────────────────
#  ROUND 1 NEW TASKS
# ─────────────────────────────────────────────────────────────

@shared_task(bind=True)
def password_protect_task(self, job_id, user_password='', owner_password=''):
    from converter.models import ConversionJob
    from converter.utils import get_output_path
    job = ConversionJob.objects.get(id=job_id)
    try:
        job.status = 'processing'; job.save()
        import fitz
        if not user_password:
            raise ValueError("A password is required to protect the PDF.")
        op = owner_password if owner_password else user_password + '_ds_owner'
        abs_path, rel_path = get_output_path(job.input_file.name, '.pdf')
        doc = fitz.open(job.input_file.path)
        permissions = (
            fitz.PDF_PERM_ACCESSIBILITY
            | fitz.PDF_PERM_PRINT
            | fitz.PDF_PERM_PRINT_HQ
        )
        doc.save(
            abs_path,
            encryption=fitz.PDF_ENCRYPT_AES_256,
            owner_pw=op,
            user_pw=user_password,
            permissions=permissions,
            garbage=3,
            deflate=True,
        )
        doc.close()
        job.output_file = rel_path
        job.output_size = os.path.getsize(abs_path)
        job.status = 'done'; job.save()
    except Exception as e:
        job.status = 'failed'; job.error_message = str(e); job.save()
        raise


@shared_task(bind=True)
def unlock_pdf_task(self, job_id, password=''):
    from converter.models import ConversionJob
    from converter.utils import get_output_path
    job = ConversionJob.objects.get(id=job_id)
    try:
        job.status = 'processing'; job.save()
        import fitz
        abs_path, rel_path = get_output_path(job.input_file.name, '.pdf')
        doc = fitz.open(job.input_file.path)
        if doc.is_encrypted:
            authenticated = doc.authenticate('')
            if not authenticated and password:
                authenticated = doc.authenticate(password)
            if not authenticated:
                raise RuntimeError("Could not unlock PDF. Check that the password is correct.")
        doc.save(abs_path, encryption=fitz.PDF_ENCRYPT_NONE, garbage=3, deflate=True)
        doc.close()
        job.output_file = rel_path
        job.output_size = os.path.getsize(abs_path)
        job.status = 'done'; job.save()
    except Exception as e:
        job.status = 'failed'; job.error_message = str(e); job.save()
        raise


@shared_task(bind=True)
def rotate_pdf_task(self, job_id, rotations=None, rotation=90, page_range='all'):
    """
    Per-page mode:  rotations={"0": 90, "2": 180}  → rotate page 0 by +90°, page 2 by +180°
    Bulk mode:      rotations=None, uses rotation + page_range (all/odd/even)
    """
    from converter.models import ConversionJob
    from converter.utils import get_output_path
    job = ConversionJob.objects.get(id=job_id)
    try:
        job.status = 'processing'; job.save()
        import fitz
        abs_path, rel_path = get_output_path(job.input_file.name, '.pdf')
        doc = fitz.open(job.input_file.path)

        if rotations:
            # Per-page rotation from live preview UI
            for page_idx_str, deg in rotations.items():
                idx = int(page_idx_str)
                if 0 <= idx < len(doc):
                    page = doc[idx]
                    current = page.rotation
                    page.set_rotation((current + int(deg)) % 360)
        else:
            # Bulk rotation from simple form
            rot = int(rotation) if rotation else 90
            if rot not in (90, 180, 270):
                rot = 90
            for i, page in enumerate(doc):
                page_num = i + 1
                should_rotate = (
                    page_range == 'all'
                    or (page_range == 'odd'  and page_num % 2 == 1)
                    or (page_range == 'even' and page_num % 2 == 0)
                )
                if should_rotate:
                    current = page.rotation
                    page.set_rotation((current + rot) % 360)

        doc.save(abs_path, garbage=3, deflate=True)
        doc.close()
        job.output_file = rel_path
        job.output_size = os.path.getsize(abs_path)
        job.status = 'done'; job.save()
    except Exception as e:
        job.status = 'failed'; job.error_message = str(e); job.save()
        raise


@shared_task(bind=True)
def watermark_pdf_task(self, job_id, watermark_text='CONFIDENTIAL',
                        opacity=30, position='diagonal'):
    from converter.models import ConversionJob
    from converter.utils import get_output_path
    job = ConversionJob.objects.get(id=job_id)
    try:
        job.status = 'processing'; job.save()
        import fitz

        text    = (watermark_text or 'CONFIDENTIAL').strip()[:50]
        opacity = max(5, min(80, int(opacity) if opacity else 30))
        # opacity % → gray: 30% opacity looks like 70% gray
        gray  = round(1.0 - (opacity / 100.0), 2)
        color = (gray, gray, gray)

        abs_path, rel_path = get_output_path(job.input_file.name, '.pdf')
        doc = fitz.open(job.input_file.path)

        for page in doc:
            rect = page.rect
            w, h = rect.width, rect.height
            font_size = max(18, min(w * 0.30 / max(len(text), 1) * 1.8, 80))
            tw = len(text) * font_size * 0.52

            if position == 'diagonal':
                # Use morph=(fixpoint, matrix) to rotate 45° — avoids bad rotate value error
                mat = fitz.Matrix(45)
                stamps = [
                    fitz.Point(w * 0.10, h * 0.35),
                    fitz.Point(w * 0.25, h * 0.60),
                    fitz.Point(w * 0.55, h * 0.25),
                    fitz.Point(w * 0.55, h * 0.75),
                    fitz.Point(w * 0.70, h * 0.50),
                ]
                for pt in stamps:
                    page.insert_text(
                        pt, text,
                        fontname='helv',
                        fontsize=font_size,
                        color=color,
                        morph=(pt, mat),
                        overlay=True,
                    )
            elif position == 'center':
                x = max(0.0, (w - tw) / 2)
                page.insert_text(fitz.Point(x, h / 2), text,
                                 fontname='helv', fontsize=font_size,
                                 color=color, overlay=True)
            elif position == 'top':
                x = max(0.0, (w - tw) / 2)
                page.insert_text(fitz.Point(x, font_size + 20), text,
                                 fontname='helv', fontsize=font_size,
                                 color=color, overlay=True)
            elif position == 'bottom':
                x = max(0.0, (w - tw) / 2)
                page.insert_text(fitz.Point(x, h - 20), text,
                                 fontname='helv', fontsize=font_size,
                                 color=color, overlay=True)

        doc.save(abs_path, garbage=3, deflate=True)
        doc.close()
        job.output_file = rel_path
        job.output_size = os.path.getsize(abs_path)
        job.status = 'done'; job.save()
    except Exception as e:
        job.status = 'failed'; job.error_message = str(e); job.save()
        raise


@shared_task(bind=True)
def add_page_numbers_task(self, job_id, position='bottom-center',
                           start_number=1, font_size=10):
    from converter.models import ConversionJob
    from converter.utils import get_output_path
    job = ConversionJob.objects.get(id=job_id)
    try:
        job.status = 'processing'; job.save()
        import fitz

        start  = max(1, int(start_number) if start_number else 1)
        fsize  = max(6, min(24, int(font_size) if font_size else 10))
        margin = 24

        abs_path, rel_path = get_output_path(job.input_file.name, '.pdf')
        doc = fitz.open(job.input_file.path)

        for i, page in enumerate(doc):
            rect   = page.rect
            label  = str(start + i)
            tw     = len(label) * fsize * 0.55
            pos_lower = position.lower()

            if 'right' in pos_lower:
                x = rect.width - tw - margin
            elif 'left' in pos_lower:
                x = margin
            else:
                x = (rect.width - tw) / 2

            if 'top' in pos_lower:
                y = margin + fsize
            else:
                y = rect.height - margin

            page.insert_text(
                fitz.Point(x, y), label,
                fontname='helv', fontsize=fsize,
                color=(0, 0, 0), overlay=True,
            )

        doc.save(abs_path, garbage=3, deflate=True)
        doc.close()
        job.output_file = rel_path
        job.output_size = os.path.getsize(abs_path)
        job.status = 'done'; job.save()
    except Exception as e:
        job.status = 'failed'; job.error_message = str(e); job.save()
        raise


"""
ROUND 2 TASKS — append these 8 functions to the END of converter/tasks.py
Do NOT replace the file — just paste everything below the last existing task.
"""

# ─────────────────────────────────────────────────────────────
#  1. PDF → Excel
#  Uses pdfplumber to extract tables page by page.
#  Each page gets its own sheet. Falls back to raw text rows
#  if no structured tables are detected.
# ─────────────────────────────────────────────────────────────

@shared_task(bind=True)
def pdf_to_excel_task(self, job_id):
    from converter.models import ConversionJob
    from converter.utils import get_output_path
    job = ConversionJob.objects.get(id=job_id)
    try:
        job.status = 'processing'; job.save()
        import pdfplumber
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

        abs_path, rel_path = get_output_path(job.input_file.name, '.xlsx')
        wb = Workbook()
        wb.remove(wb.active)   # remove default empty sheet

        header_font  = Font(bold=True, color='FFFFFF', size=10)
        header_fill  = PatternFill('solid', fgColor='1a73e8')
        center_align = Alignment(horizontal='center', vertical='center', wrap_text=True)
        left_align   = Alignment(horizontal='left',   vertical='center', wrap_text=True)
        thin_border  = Border(
            left=Side(style='thin', color='DADCE0'),
            right=Side(style='thin', color='DADCE0'),
            top=Side(style='thin', color='DADCE0'),
            bottom=Side(style='thin', color='DADCE0'),
        )

        with pdfplumber.open(job.input_file.path) as pdf:
            for page_num, page in enumerate(pdf.pages, start=1):
                sheet_name = f'Page {page_num}'
                ws = wb.create_sheet(title=sheet_name)
                ws.row_dimensions[1].height = 20

                tables = page.extract_tables()

                if tables:
                    # ── Structured tables found ───────────────
                    row_cursor = 1
                    for table_idx, table in enumerate(tables):
                        if table_idx > 0:
                            row_cursor += 1  # blank row between tables

                        for r_idx, row in enumerate(table):
                            for c_idx, cell in enumerate(row):
                                cell_val = str(cell).strip() if cell else ''
                                ws_cell  = ws.cell(row=row_cursor, column=c_idx + 1, value=cell_val)
                                ws_cell.border    = thin_border
                                ws_cell.alignment = left_align
                                if r_idx == 0:
                                    # Header row styling
                                    ws_cell.font      = header_font
                                    ws_cell.fill      = header_fill
                                    ws_cell.alignment = center_align
                            row_cursor += 1
                else:
                    # ── No tables — fall back to raw text lines ──
                    text = page.extract_text() or ''
                    lines = [l.strip() for l in text.split('\n') if l.strip()]
                    if not lines:
                        ws.cell(row=1, column=1, value=f'(Page {page_num} — no text found)')
                    else:
                        ws.cell(row=1, column=1, value=f'Page {page_num} — extracted text')
                        ws.cell(row=1, column=1).font = header_font
                        ws.cell(row=1, column=1).fill = header_fill
                        for r_idx, line in enumerate(lines, start=2):
                            ws_cell = ws.cell(row=r_idx, column=1, value=line)
                            ws_cell.alignment = left_align

                # Auto-size columns (cap at 60)
                for col in ws.columns:
                    max_len = 0
                    col_letter = col[0].column_letter
                    for cell in col:
                        try:
                            if cell.value:
                                max_len = max(max_len, len(str(cell.value)))
                        except Exception:
                            pass
                    ws.column_dimensions[col_letter].width = min(max(max_len + 2, 10), 60)

        if not wb.sheetnames:
            ws = wb.create_sheet('Sheet1')
            ws.cell(row=1, column=1, value='No content found in PDF')

        wb.save(abs_path)
        job.output_file = rel_path
        job.output_size = os.path.getsize(abs_path)
        job.status = 'done'; job.save()
    except Exception as e:
        job.status = 'failed'; job.error_message = str(e); job.save()
        raise


# ─────────────────────────────────────────────────────────────
#  2. Excel → PDF
#  Reads each sheet with openpyxl, writes a formatted PDF
#  using reportlab Table objects so column widths are preserved.
# ─────────────────────────────────────────────────────────────

@shared_task(bind=True)
def excel_to_pdf_task(self, job_id):
    from converter.models import ConversionJob
    from converter.utils import get_output_path
    job = ConversionJob.objects.get(id=job_id)
    try:
        job.status = 'processing'; job.save()
        from openpyxl import load_workbook
        from reportlab.lib.pagesizes import A4, landscape
        from reportlab.lib.units import inch, cm
        from reportlab.lib import colors
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.platypus import (SimpleDocTemplate, Table, TableStyle,
                                         Paragraph, Spacer, PageBreak)
        from reportlab.lib.enums import TA_CENTER

        abs_path, rel_path = get_output_path(job.input_file.name, '.pdf')
        wb = load_workbook(job.input_file.path, data_only=True)

        styles   = getSampleStyleSheet()
        title_st = ParagraphStyle('Title', parent=styles['Heading1'],
                                   fontSize=13, spaceAfter=8, alignment=TA_CENTER)
        cell_st  = ParagraphStyle('Cell',  parent=styles['Normal'],
                                   fontSize=8, leading=10)

        story = []

        for sheet_idx, ws in enumerate(wb.worksheets):
            if sheet_idx > 0:
                story.append(PageBreak())

            story.append(Paragraph(ws.title, title_st))
            story.append(Spacer(1, 0.15 * inch))

            # Collect rows — skip fully empty rows
            rows_data = []
            for row in ws.iter_rows(values_only=True):
                if any(cell is not None for cell in row):
                    rows_data.append([
                        Paragraph(str(c) if c is not None else '', cell_st)
                        for c in row
                    ])

            if not rows_data:
                story.append(Paragraph('(Empty sheet)', cell_st))
                continue

            # Page width minus margins, split evenly across columns
            col_count  = max(len(r) for r in rows_data)
            page_w     = A4[0] - 1.4 * inch   # usable width in points
            col_w      = page_w / max(col_count, 1)
            col_widths = [col_w] * col_count

            tbl = Table(rows_data, colWidths=col_widths, repeatRows=1)
            tbl.setStyle(TableStyle([
                # Header row (row 0)
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#1a73e8')),
                ('TEXTCOLOR',  (0, 0), (-1, 0), colors.white),
                ('FONTNAME',   (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE',   (0, 0), (-1, 0), 9),
                ('ALIGN',      (0, 0), (-1, 0), 'CENTER'),
                # Data rows
                ('FONTNAME',   (0, 1), (-1, -1), 'Helvetica'),
                ('FONTSIZE',   (0, 1), (-1, -1), 8),
                ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#f8f9fa')]),
                ('GRID',       (0, 0), (-1, -1), 0.4, colors.HexColor('#dadce0')),
                ('VALIGN',     (0, 0), (-1, -1), 'TOP'),
                ('LEFTPADDING',  (0, 0), (-1, -1), 4),
                ('RIGHTPADDING', (0, 0), (-1, -1), 4),
                ('TOPPADDING',   (0, 0), (-1, -1), 3),
                ('BOTTOMPADDING',(0, 0), (-1, -1), 3),
            ]))
            story.append(tbl)

        pdf = SimpleDocTemplate(abs_path, pagesize=A4,
                                 rightMargin=0.7 * inch, leftMargin=0.7 * inch,
                                 topMargin=0.7 * inch,   bottomMargin=0.7 * inch)
        pdf.build(story)
        job.output_file = rel_path
        job.output_size = os.path.getsize(abs_path)
        job.status = 'done'; job.save()
    except Exception as e:
        job.status = 'failed'; job.error_message = str(e); job.save()
        raise


# ─────────────────────────────────────────────────────────────
#  3. PowerPoint → PDF
#  Renders each slide as a full-page image using python-pptx
#  to extract content, then writes each slide's text/shapes
#  as a reportlab page. For maximum reliability we render
#  the slide as a thumbnail image via python-pptx + Pillow.
#
#  NOTE: python-pptx can't render slides to images directly.
#  We extract text and render it as structured PDF pages.
#  For pixel-perfect rendering, LibreOffice would be needed,
#  but that's a system dependency we avoid.
# ─────────────────────────────────────────────────────────────

@shared_task(bind=True)
def pptx_to_pdf_task(self, job_id):
    from converter.models import ConversionJob
    from converter.utils import get_output_path
    job = ConversionJob.objects.get(id=job_id)
    try:
        job.status = 'processing'; job.save()
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from reportlab.lib.pagesizes import landscape, A4
        from reportlab.lib.units import inch
        from reportlab.lib import colors
        from reportlab.lib.styles import ParagraphStyle
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
        from reportlab.lib.enums import TA_LEFT, TA_CENTER

        abs_path, rel_path = get_output_path(job.input_file.name, '.pdf')
        prs = Presentation(job.input_file.path)

        # Slide dimensions (default 10×7.5 inches in pptx)
        PAGE_W, PAGE_H = landscape(A4)

        title_st = ParagraphStyle('SlideTitle', fontName='Helvetica-Bold',
                                   fontSize=20, textColor=colors.HexColor('#202124'),
                                   spaceAfter=12, leading=26)
        body_st  = ParagraphStyle('SlideBody',  fontName='Helvetica',
                                   fontSize=12, textColor=colors.HexColor('#3c4043'),
                                   spaceAfter=6, leading=16, leftIndent=12)
        note_st  = ParagraphStyle('SlideNum',   fontName='Helvetica',
                                   fontSize=8, textColor=colors.HexColor('#9aa0a6'),
                                   spaceAfter=4)
        bullet_st = ParagraphStyle('Bullet',    fontName='Helvetica',
                                    fontSize=11, textColor=colors.HexColor('#3c4043'),
                                    spaceAfter=4, leading=14, leftIndent=20,
                                    bulletIndent=8, bulletText='•')

        story = []

        for slide_idx, slide in enumerate(prs.slides):
            if slide_idx > 0:
                story.append(PageBreak())

            story.append(Paragraph(f'Slide {slide_idx + 1}', note_st))

            # Extract text from all shapes in reading order (top→bottom)
            shapes_with_top = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    top = shape.top or 0
                    shapes_with_top.append((top, shape))
            shapes_with_top.sort(key=lambda x: x[0])

            slide_has_content = False
            for _, shape in shapes_with_top:
                tf = shape.text_frame
                for para_idx, para in enumerate(tf.paragraphs):
                    text = para.text.strip()
                    if not text:
                        continue
                    slide_has_content = True

                    # Detect title heuristic: first shape, first paragraph, larger font
                    is_title = (
                        shape == shapes_with_top[0][1]
                        and para_idx == 0
                        and any(run.font.size and run.font.size >= Pt(18)
                                for run in para.runs if run.font.size)
                    ) if shapes_with_top else False

                    # Use bold style for title-like text
                    if is_title or (shape == shapes_with_top[0][1] and para_idx == 0):
                        safe = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                        story.append(Paragraph(safe, title_st))
                    else:
                        safe = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                        story.append(Paragraph(safe, bullet_st if para.level > 0 else body_st))

            if not slide_has_content:
                story.append(Paragraph('(No text content)', note_st))

            story.append(Spacer(1, 0.2 * inch))

        if not story:
            story.append(Paragraph('(Empty presentation)', body_st))

        pdf = SimpleDocTemplate(abs_path, pagesize=landscape(A4),
                                 rightMargin=0.8 * inch, leftMargin=0.8 * inch,
                                 topMargin=0.6 * inch,   bottomMargin=0.6 * inch)
        pdf.build(story)
        job.output_file = rel_path
        job.output_size = os.path.getsize(abs_path)
        job.status = 'done'; job.save()
    except Exception as e:
        job.status = 'failed'; job.error_message = str(e); job.save()
        raise


# ─────────────────────────────────────────────────────────────
#  4. PDF → PowerPoint
#  Renders each PDF page as a full-slide image (PNG via fitz),
#  then creates a PPTX with one image-per-slide.
#  This approach preserves the visual layout perfectly.
# ─────────────────────────────────────────────────────────────

@shared_task(bind=True)
def pdf_to_pptx_task(self, job_id):
    from converter.models import ConversionJob
    from converter.utils import get_output_path
    job = ConversionJob.objects.get(id=job_id)
    try:
        job.status = 'processing'; job.save()
        import fitz
        from pptx import Presentation
        from pptx.util import Inches, Emu
        import io

        abs_path, rel_path = get_output_path(job.input_file.name, '.pptx')
        doc = fitz.open(job.input_file.path)

        prs = Presentation()

        # Remove default blank slide layout content
        blank_layout = prs.slide_layouts[6]  # completely blank layout

        for page_num, page in enumerate(doc):
            # Render page at 150 DPI
            mat = fitz.Matrix(150 / 72, 150 / 72)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            img_bytes = pix.tobytes('png')

            # Set slide dimensions to match PDF page aspect ratio
            page_w = page.rect.width   # in points
            page_h = page.rect.height  # in points

            # Widescreen default is 10×5.625 inches — override to match PDF
            prs.slide_width  = Inches(10)
            prs.slide_height = Inches(10 * page_h / page_w)

            slide = prs.slides.add_slide(blank_layout)

            # Add image covering the entire slide
            pic_stream = io.BytesIO(img_bytes)
            slide.shapes.add_picture(
                pic_stream,
                left=Inches(0),
                top=Inches(0),
                width=prs.slide_width,
                height=prs.slide_height,
            )

        doc.close()
        prs.save(abs_path)
        job.output_file = rel_path
        job.output_size = os.path.getsize(abs_path)
        job.status = 'done'; job.save()
    except Exception as e:
        job.status = 'failed'; job.error_message = str(e); job.save()
        raise


# ─────────────────────────────────────────────────────────────
#  6. OCR — scanned PDF → searchable PDF
#  Renders each page as an image, runs Tesseract OCR,
#  then uses fitz to embed an invisible text layer over the
#  original image so the PDF becomes copy-pasteable.
# ─────────────────────────────────────────────────────────────
"""
FILE: two_fixed_tasks.py
========================
These are REPLACEMENT functions for your converter/tasks.py

For ocr_pdf_task:
  Find the existing @shared_task def ocr_pdf_task(...) in your tasks.py
  and replace the ENTIRE function with the one below.

For html_to_pdf_task:
  Find the existing @shared_task def html_to_pdf_task(...) in your tasks.py
  and replace the ENTIRE function with the one below.
"""



@shared_task(bind=True)
def ocr_pdf_task(self, job_id):
    from converter.models import ConversionJob
    from converter.utils import get_output_path
    job = ConversionJob.objects.get(id=job_id)
    try:
        job.status = 'processing'; job.save()
        import fitz
        import pytesseract
        from PIL import Image
        import io
        from django.conf import settings as django_settings

        # ── Fix: set both the exe path AND tessdata prefix ──────────────
        if hasattr(django_settings, 'TESSERACT_CMD'):
            pytesseract.pytesseract.tesseract_cmd = django_settings.TESSERACT_CMD

        # Set TESSDATA_PREFIX to the folder containing the .traineddata files
        # This fixes "Error opening data file eng.traineddata" on Windows
        tesseract_dir = os.path.dirname(
            getattr(django_settings, 'TESSERACT_CMD',
                    r'C:\Program Files\Tesseract-OCR\tesseract.exe')
        )
        os.environ['TESSDATA_PREFIX'] = os.path.join(tesseract_dir, 'tessdata')
        # ────────────────────────────────────────────────────────────────

        abs_path, rel_path = get_output_path(job.input_file.name, '.pdf')

        src_doc = fitz.open(job.input_file.path)
        out_doc = fitz.open()

        for page_num, page in enumerate(src_doc):
            page_rect = page.rect
            w, h = page_rect.width, page_rect.height

            # Render at 200 DPI for good OCR accuracy
            mat = fitz.Matrix(200 / 72, 200 / 72)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            img_bytes = pix.tobytes('png')

            # Run Tesseract
            pil_img  = Image.open(io.BytesIO(img_bytes))
            ocr_data = pytesseract.image_to_data(
                pil_img,
                output_type=pytesseract.Output.DICT,
                config='--psm 3'
            )

            # New page, same size as original
            new_page = out_doc.new_page(width=w, height=h)

            # Background image
            new_page.insert_image(fitz.Rect(0, 0, w, h), stream=img_bytes)

            # Scale factors: image pixels → PDF points
            scale_x = w / pix.width
            scale_y = h / pix.height

            # Invisible text layer for searchability
            n = len(ocr_data['text'])
            for i in range(n):
                word = (ocr_data['text'][i] or '').strip()
                conf = int(ocr_data['conf'][i])
                if not word or conf < 30:
                    continue

                ix = ocr_data['left'][i]  * scale_x
                iy = ocr_data['top'][i]   * scale_y
                iw = ocr_data['width'][i] * scale_x
                ih = ocr_data['height'][i] * scale_y

                if iw <= 0 or ih <= 0:
                    continue

                font_size = max(4.0, ih * 0.72)
                try:
                    new_page.insert_text(
                        fitz.Point(ix, iy + ih),
                        word,
                        fontname='helv',
                        fontsize=font_size,
                        color=(1, 1, 1),   # white = invisible over image
                        overlay=False,
                    )
                except Exception:
                    pass

        src_doc.close()
        out_doc.save(abs_path, garbage=3, deflate=True)
        out_doc.close()

        job.output_file = rel_path
        job.output_size = os.path.getsize(abs_path)
        job.status = 'done'; job.save()
    except Exception as e:
        job.status = 'failed'; job.error_message = str(e); job.save()
        raise


@shared_task(bind=True)
def html_to_pdf_task(self, job_id):
    from converter.models import ConversionJob
    from converter.utils import get_output_path
    job = ConversionJob.objects.get(id=job_id)
    try:
        job.status = 'processing'; job.save()
        import re
        from xhtml2pdf import pisa

        abs_path, rel_path = get_output_path(job.input_file.name, '.pdf')

        with open(job.input_file.path, 'r', encoding='utf-8', errors='replace') as f:
            html_content = f.read()

        # ── Fix: strip external resources that crash xhtml2pdf on Windows ──
        # Remove <link rel="stylesheet"> (Google Fonts, CDN stylesheets)
        html_content = re.sub(
            r'<link[^>]+rel=["\']stylesheet["\'][^>]*>',
            '', html_content, flags=re.IGNORECASE)
        # Remove @import rules (e.g. @import url('https://fonts.googleapis.com/...'))
        html_content = re.sub(
            r'@import\s+[^;]+;',
            '', html_content, flags=re.IGNORECASE)
        # Remove @font-face blocks (tries to load TTF files from temp — fails on Windows)
        html_content = re.sub(
            r'@font-face\s*\{[^}]*\}',
            '', html_content, flags=re.IGNORECASE | re.DOTALL)
        # ────────────────────────────────────────────────────────────────────

        # Inject safe base CSS using only built-in PDF fonts (Helvetica, Courier)
        base_css = """
        <style>
          body { font-family: Helvetica, Arial, sans-serif; font-size: 11pt;
                 line-height: 1.5; color: #202124; margin: 0; padding: 0; }
          h1 { font-size: 20pt; font-weight: bold; margin: 14pt 0 6pt; }
          h2 { font-size: 16pt; font-weight: bold; margin: 12pt 0 5pt; }
          h3 { font-size: 13pt; font-weight: bold; margin: 10pt 0 4pt; }
          h4 { font-size: 11pt; font-weight: bold; margin: 8pt 0 3pt; }
          p  { margin: 0 0 6pt; }
          table { border-collapse: collapse; width: 100%; margin: 8pt 0; }
          th, td { border: 1px solid #dadce0; padding: 4pt 6pt; font-size: 10pt; }
          th { background: #1a73e8; color: #ffffff; font-weight: bold; }
          tr:nth-child(even) { background: #f8f9fa; }
          pre, code { font-family: Courier, monospace; font-size: 9pt;
                      background: #f8f9fa; padding: 2pt 4pt; }
          blockquote { border-left: 3pt solid #1a73e8; margin-left: 12pt;
                       padding-left: 8pt; color: #5f6368; }
          a { color: #1a73e8; }
          img { max-width: 100%; }
          ul, ol { margin: 4pt 0 4pt 16pt; padding: 0; }
          li { margin-bottom: 2pt; }
          hr { border: none; border-top: 1pt solid #dadce0; margin: 10pt 0; }
        </style>
        """

        if '</head>' in html_content.lower():
            idx = html_content.lower().index('</head>')
            html_content = html_content[:idx] + base_css + html_content[idx:]
        else:
            html_content = (
                '<html><head>' + base_css + '</head><body>'
                + html_content
                + '</body></html>'
            )

        with open(abs_path, 'wb') as out_f:
            result = pisa.CreatePDF(html_content, dest=out_f)

        if result.err:
            raise RuntimeError(f"HTML to PDF conversion error (code {result.err}). "
                               "Check that the HTML is valid and contains no external resources.")

        job.output_file = rel_path
        job.output_size = os.path.getsize(abs_path)
        job.status = 'done'; job.save()
    except Exception as e:
        job.status = 'failed'; job.error_message = str(e); job.save()
        raise


# ─────────────────────────────────────────────────────────────
#  7. Extract Text from PDF → .txt
#  Uses fitz to extract all text, page by page.
#  Outputs a clean .txt file with page separators.
# ─────────────────────────────────────────────────────────────

@shared_task(bind=True)
def extract_text_task(self, job_id):
    from converter.models import ConversionJob
    from converter.utils import get_output_path
    job = ConversionJob.objects.get(id=job_id)
    try:
        job.status = 'processing'; job.save()
        import fitz

        abs_path, rel_path = get_output_path(job.input_file.name, '.txt')
        doc   = fitz.open(job.input_file.path)
        lines = []

        for page_num, page in enumerate(doc, start=1):
            text = page.get_text('text').strip()
            lines.append(f'{'─' * 60}')
            lines.append(f'  Page {page_num} of {len(doc)}')
            lines.append(f'{'─' * 60}')
            if text:
                lines.append(text)
            else:
                lines.append('(No text found on this page — may be a scanned image)')
            lines.append('')

        doc.close()

        full_text = '\n'.join(lines)
        with open(abs_path, 'w', encoding='utf-8') as f:
            f.write(full_text)

        job.output_file = rel_path
        job.output_size = os.path.getsize(abs_path)
        job.status = 'done'; job.save()
    except Exception as e:
        job.status = 'failed'; job.error_message = str(e); job.save()
        raise


# ─────────────────────────────────────────────────────────────
#  8. Extract Images from PDF → .zip
#  Uses fitz to pull every embedded image from the PDF.
#  Saves each as PNG in a zip file.
#  Skips tiny images (< 10×10 px) which are usually decorative.
# ─────────────────────────────────────────────────────────────

@shared_task(bind=True)
def extract_images_task(self, job_id):
    from converter.models import ConversionJob
    from converter.utils import get_output_path, create_zip
    job = ConversionJob.objects.get(id=job_id)
    try:
        job.status = 'processing'; job.save()
        import fitz
        from PIL import Image
        import io

        doc      = fitz.open(job.input_file.path)
        tmp_dir  = tempfile.mkdtemp()
        files_dict  = {}
        img_counter = 0
        MIN_SIZE    = 10   # skip images smaller than 10×10 px

        for page_num, page in enumerate(doc, start=1):
            image_list = page.get_images(full=True)
            for img_info in image_list:
                xref = img_info[0]
                try:
                    base_image = doc.extract_image(xref)
                    img_bytes  = base_image['image']
                    img_ext    = base_image['ext']   # 'png', 'jpeg', etc.

                    # Check dimensions via PIL — skip tiny decorative images
                    pil_img = Image.open(io.BytesIO(img_bytes))
                    iw, ih  = pil_img.size
                    if iw < MIN_SIZE or ih < MIN_SIZE:
                        continue

                    img_counter += 1
                    # Always save as PNG for consistency
                    if pil_img.mode in ('RGBA', 'P', 'LA'):
                        pil_img = pil_img.convert('RGBA')
                    elif pil_img.mode not in ('RGB', 'RGBA', 'L'):
                        pil_img = pil_img.convert('RGB')

                    filename  = f'page{page_num:03d}_img{img_counter:03d}.png'
                    save_path = os.path.join(tmp_dir, filename)
                    pil_img.save(save_path, 'PNG', optimize=True)
                    files_dict[filename] = save_path

                except Exception:
                    continue   # skip unreadable images silently

        doc.close()

        if not files_dict:
            raise RuntimeError(
                'No images found in this PDF. '
                'The file may contain no embedded images, or only very small decorative elements.'
            )

        zip_abs, zip_rel = create_zip(files_dict, 'extracted_images.zip')
        for fp in files_dict.values():
            try: os.remove(fp)
            except: pass

        job.output_file = zip_rel
        job.output_size = os.path.getsize(zip_abs)
        job.status = 'done'; job.save()
    except Exception as e:
        job.status = 'failed'; job.error_message = str(e); job.save()
        raise


# ─────────────────────────────────────────────────────────────
#  Round 6 — Tool Enhancements
# ─────────────────────────────────────────────────────────────

@shared_task(bind=True)
def edit_metadata_task(self, job_id, title='', author='', subject='', keywords=''):
    from converter.models import ConversionJob
    from converter.utils import get_output_path
    job = ConversionJob.objects.get(id=job_id)
    try:
        job.status = 'processing'; job.save()
        import fitz
        doc = fitz.open(job.input_file.path)
        meta = doc.metadata
        if title: meta['title'] = title
        if author: meta['author'] = author
        if subject: meta['subject'] = subject
        if keywords: meta['keywords'] = keywords
        doc.set_metadata(meta)
        abs_path, rel_path = get_output_path(job.input_file.name, '.pdf')
        doc.save(abs_path)
        doc.close()
        job.output_file = rel_path
        job.output_size = os.path.getsize(abs_path)
        job.status = 'done'; job.save()
    except Exception as e:
        job.status = 'failed'; job.error_message = str(e); job.save()
        raise


@shared_task(bind=True)
def flatten_pdf_task(self, job_id):
    from converter.models import ConversionJob
    from converter.utils import get_output_path
    job = ConversionJob.objects.get(id=job_id)
    try:
        job.status = 'processing'; job.save()
        import fitz
        doc = fitz.open(job.input_file.path)
        out_doc = fitz.open()
        for page in doc:
            rect = page.rect
            pix = page.get_pixmap(dpi=150)
            out_page = out_doc.new_page(width=rect.width, height=rect.height)
            out_page.insert_image(rect, pixmap=pix)
        abs_path, rel_path = get_output_path(job.input_file.name, '.pdf')
        out_doc.save(abs_path, deflate=True)
        out_doc.close()
        doc.close()
        job.output_file = rel_path
        job.output_size = os.path.getsize(abs_path)
        job.status = 'done'; job.save()
    except Exception as e:
        job.status = 'failed'; job.error_message = str(e); job.save()
        raise


@shared_task(bind=True)
def grayscale_pdf_task(self, job_id):
    from converter.models import ConversionJob
    from converter.utils import get_output_path
    job = ConversionJob.objects.get(id=job_id)
    try:
        job.status = 'processing'; job.save()
        import fitz
        doc = fitz.open(job.input_file.path)
        out_doc = fitz.open()
        for page in doc:
            rect = page.rect
            pix = page.get_pixmap(dpi=150, colorspace=fitz.csGRAY)
            out_page = out_doc.new_page(width=rect.width, height=rect.height)
            out_page.insert_image(rect, pixmap=pix)
        abs_path, rel_path = get_output_path(job.input_file.name, '.pdf')
        out_doc.save(abs_path, deflate=True)
        out_doc.close()
        doc.close()
        job.output_file = rel_path
        job.output_size = os.path.getsize(abs_path)
        job.status = 'done'; job.save()
    except Exception as e:
        job.status = 'failed'; job.error_message = str(e); job.save()
        raise


@shared_task(bind=True)
def crop_pdf_task(self, job_id, margin_top=0, margin_bottom=0, margin_left=0, margin_right=0):
    from converter.models import ConversionJob
    from converter.utils import get_output_path
    job = ConversionJob.objects.get(id=job_id)
    try:
        job.status = 'processing'; job.save()
        import fitz
        doc = fitz.open(job.input_file.path)
        for page in doc:
            r = page.rect
            page.set_cropbox(fitz.Rect(
                r.x0 + margin_left,
                r.y0 + margin_top,
                r.x1 - margin_right,
                r.y1 - margin_bottom
            ))
        abs_path, rel_path = get_output_path(job.input_file.name, '.pdf')
        doc.save(abs_path)
        doc.close()
        job.output_file = rel_path
        job.output_size = os.path.getsize(abs_path)
        job.status = 'done'; job.save()
    except Exception as e:
        job.status = 'failed'; job.error_message = str(e); job.save()
        raise


@shared_task(bind=True)
def render_rotate_thumbnails_task(self, job_id):
    from converter.models import ConversionJob
    from django.core.files.storage import default_storage
    from django.core.files.base import ContentFile
    import fitz
    job = ConversionJob.objects.get(id=job_id)
    try:
        job.status = 'analysing'
        job.save(update_fields=['status'])
        
        # doc.path will trigger S3 download via s3_patch.py on the WORKER (Railway)
        # This is fine on Railway, as it has more memory and no gunicorn timeout.
        doc = fitz.open(job.input_file.path)
        mat = fitz.Matrix(96 / 72, 96 / 72)
        for i, page in enumerate(doc):
            pix = page.get_pixmap(matrix=mat, alpha=False)
            path = f"rotate_previews/{job.id}/page_{i}.jpg"
            default_storage.save(path, ContentFile(pix.tobytes("jpeg")))
            
        page_count = len(doc)
        doc.close()
        
        # Mark as done so the UI knows to redirect
        job.status = 'done' 
        job.error_message = f'PAGES:{page_count}'
        job.save(update_fields=['status', 'error_message'])
    except Exception as e:
        job.status = 'failed'
        job.error_message = str(e)
        job.save(update_fields=['status', 'error_message'])
        raise
